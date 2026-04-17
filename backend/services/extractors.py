from __future__ import annotations

from pathlib import Path
from typing import Literal

from docx import Document as DocxDocument
from pypdf import PdfReader
from pptx import Presentation

ContentType = Literal["text", "markdown", "pdf", "docx", "pptx", "image"]

# Map file extensions to content types. Lowercase comparisons only.
EXTENSION_MAP: dict[str, ContentType] = {
    ".txt": "text",
    ".md": "markdown",
    ".markdown": "markdown",
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".gif": "image",
    ".webp": "image",
    ".bmp": "image",
}


def detect_content_type(path: str | Path) -> ContentType:
    suffix = Path(path).suffix.lower()
    if suffix not in EXTENSION_MAP:
        raise ValueError(f"Unsupported file extension: {suffix}")
    return EXTENSION_MAP[suffix]


def read_text_file(path: str | Path) -> str:
    return Path(path).read_text(encoding="utf-8", errors="replace")


def extract_text_from_pdf(path: str | Path) -> str:
    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text)
    return "\n\n".join(parts)


def extract_text_from_docx(path: str | Path) -> str:
    doc = DocxDocument(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_from_pptx(path: str | Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def extract_body(path: str | Path, content_type: ContentType) -> str:
    """Return extracted text for indexing. Empty string for images (no OCR in v1)."""
    if content_type in ("text", "markdown"):
        return read_text_file(path)
    if content_type == "pdf":
        return extract_text_from_pdf(path)
    if content_type == "docx":
        return extract_text_from_docx(path)
    if content_type == "pptx":
        return extract_text_from_pptx(path)
    if content_type == "image":
        return ""
    raise ValueError(f"Unknown content type: {content_type}")
